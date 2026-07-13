#!/usr/bin/env python3
"""Generate DOCX / XLSX / PPTX fixtures with KNOWN content + ground-truth probes,
to substantiate Docling's 'unified multi-format' claim beyond HTML/PDF.
Run with the docling venv python (has python-docx / openpyxl / python-pptx)."""
import json, os
OUT = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "artifacts", "fixtures", "docs"))
os.makedirs(OUT, exist_ok=True)
gt = {}

# ---- DOCX: headings + a merged-cell table + a bullet list ----
from docx import Document
d = Document()
d.add_heading("Quarterly Operations Review", level=1)
d.add_heading("Summary", level=2)
d.add_paragraph("Revenue grew across every region in Q2 2026.")
d.add_heading("Regional detail", level=2)
t = d.add_table(rows=4, cols=3)
t.style = "Table Grid"
cells = [["Region", "Units", "Value"],
         ["EMEA", "1310", "19650"],
         ["APAC", "970", "14550"],
         ["AMER", "2510", "50200"]]
for r in range(4):
    for c in range(3):
        t.cell(r, c).text = cells[r][c]
# merge first two cells of a new bottom row into a 'Total' label spanning 2 cols
row = t.add_row().cells
row[0].text = "Total"
row[0].merge(row[1])
row[2].text = "84400"
d.add_heading("Action items", level=2)
for item in ["Expand APAC warehouse", "Renegotiate EMEA freight", "Pilot AMER same-day"]:
    d.add_paragraph(item, style="List Bullet")
docx_path = os.path.join(OUT, "report.docx")
d.save(docx_path)
gt["report.docx"] = {"probes": ["Quarterly Operations Review", "Regional detail", "EMEA", "50200",
                                 "Total", "84400", "Expand APAC warehouse"],
                     "expects_table": True, "expects_merged_total_row": True,
                     "headings": {"h1": 1, "h2": 3}}
print("docx:", os.path.getsize(docx_path), "bytes")

# ---- XLSX: two sheets, formulas evaluate to values, a blank column ----
import openpyxl
wb = openpyxl.Workbook()
ws = wb.active
ws.title = "Sales"
for row in [["Item", "Jan", "Blank", "Feb", "Total"],
            ["Alpha", 100, None, 120, 220],
            ["Beta", 210, None, 231, 441],
            ["Gamma", 55, None, 62, 117]]:
    ws.append(row)
ws2 = wb.create_sheet("Meta")
ws2.append(["Key", "Value"])
ws2.append(["prepared_by", "ops"])
ws2.append(["fiscal_year", 2026])
xlsx_path = os.path.join(OUT, "workbook.xlsx")
wb.save(xlsx_path)
gt["workbook.xlsx"] = {"probes": ["Alpha", "220", "Beta", "441", "prepared_by", "2026"],
                       "sheets": ["Sales", "Meta"], "expects_blank_col": True}
print("xlsx:", os.path.getsize(xlsx_path), "bytes")

# ---- PPTX: 3 slides with titles + a bullet slide + a table slide ----
from pptx import Presentation
from pptx.util import Inches
prs = Presentation()
s1 = prs.slides.add_slide(prs.slide_layouts[0])
s1.shapes.title.text = "Docling Multi-Format Test Deck"
s1.placeholders[1].text = "Slide subtitle goes here"
s2 = prs.slides.add_slide(prs.slide_layouts[1])
s2.shapes.title.text = "Key Points"
tf = s2.placeholders[1].text_frame
tf.text = "First bullet about extraction"
for extra in ["Second bullet about tables", "Third bullet about formats"]:
    tf.add_paragraph().text = extra
s3 = prs.slides.add_slide(prs.slide_layouts[5])
s3.shapes.title.text = "Numbers"
rows, cols = 3, 3
gtbl = s3.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(6), Inches(2)).table
tvals = [["Metric", "Q1", "Q2"], ["Revenue", "1204", "1388"], ["Margin", "50%", "52%"]]
for r in range(rows):
    for c in range(cols):
        gtbl.cell(r, c).text = tvals[r][c]
pptx_path = os.path.join(OUT, "deck.pptx")
prs.save(pptx_path)
gt["deck.pptx"] = {"probes": ["Docling Multi-Format Test Deck", "Key Points",
                              "Second bullet about tables", "Revenue", "1388", "Margin"],
                   "slides": 3, "expects_table": True}
print("pptx:", os.path.getsize(pptx_path), "bytes")

with open(os.path.join(OUT, "DOC_GROUNDTRUTH.json"), "w") as f:
    json.dump(gt, f, indent=2)
print("ground truth written; fixtures in", OUT)
